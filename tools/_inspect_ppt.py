from pptx import Presentation
from pptx.util import Inches, Emu

prs = Presentation(r"docs\FAE_Sharing_Intelligence_Platform_Project_v2_EN.pptx")
for si, s in enumerate(prs.slides, 1):
    print("SLIDE", si)
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.split("\n")[0][:50]
            print(
                " ",
                sh.name,
                "l",
                round(sh.left.inches, 2),
                "t",
                round(sh.top.inches, 2),
                "w",
                round(sh.width.inches, 2),
                "h",
                round(sh.height.inches, 2),
                "|",
                t,
            )
