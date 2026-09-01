from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

p = r"docs\FAE_Sharing_Intelligence_Platform_Project_v2_EN.pptx"
prs = Presentation(p)
print("slides", len(prs.slides), "w", prs.slide_width, "h", prs.slide_height)
print("layouts", [l.name for l in prs.slide_layouts])


def dump_shape(sh, indent="  "):
    kind = sh.shape_type
    name = sh.name
    if sh.has_table:
        tbl = sh.table
        print(indent + "[TABLE] %dx%d %s" % (len(tbl.rows), len(tbl.columns), name))
        for row in tbl.rows:
            cells = [c.text.replace("\n", " ")[:60] for c in row.cells]
            print(indent + "  " + str(cells))
        return
    if sh.has_text_frame:
        t = sh.text_frame.text.replace("\n", " | ")
        if t.strip():
            print(indent + "[T] " + name + ": " + t[:300])
        return
    if kind == MSO_SHAPE_TYPE.PICTURE:
        print(indent + "[PIC] " + name)
        return
    if kind == MSO_SHAPE_TYPE.GROUP:
        print(indent + "[GROUP] " + name)
        for g in sh.shapes:
            dump_shape(g, indent + "  ")


for i, s in enumerate(prs.slides, 1):
    print("=" * 60)
    print("SLIDE", i)
    for sh in s.shapes:
        dump_shape(sh)
