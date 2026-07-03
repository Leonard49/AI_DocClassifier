"""PPT (.pptx) attachment extractor."""

from __future__ import annotations

from .base import BaseExtractor


class PPTExtractor(BaseExtractor):
    """Extract text and images from PowerPoint attachments into docx blocks."""

    def extract(self, pptx_path: str, doc_token: str, root_block_id: str) -> None:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            print("    缺少 python-pptx")
            self.append_blocks(
                doc_token, [self._placeholder("请安装: pip install python-pptx")]
            )
            return

        prs = Presentation(pptx_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            items = []
            for shape in slide.shapes:
                top = shape.top if shape.top else 0
                left = shape.left if shape.left else 0

                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        items.append((top, left, "text", text))
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    items.append(
                        (
                            top,
                            left,
                            "image",
                            {
                                "image_bytes": img.blob,
                                "image_ext": img.content_type.split("/")[-1],
                            },
                        )
                    )

            if not items:
                continue

            items.sort(key=lambda x: (x[0], x[1]))
            text_buf = []

            for _top, _left, item_type, data in items:
                if item_type == "text":
                    clean = self.sanitize(data)
                    if clean:
                        text_buf.append(
                            {
                                "block_type": 2,
                                "text": {
                                    "elements": [{"text_run": {"content": clean}}]
                                },
                            }
                        )
                elif item_type == "image":
                    if text_buf:
                        self.append_blocks(doc_token, text_buf)
                        text_buf = []
                    self.insert_image(
                        doc_token,
                        root_block_id,
                        data["image_bytes"],
                        data["image_ext"],
                        len(data["image_bytes"]),
                    )

            if text_buf:
                self.append_blocks(doc_token, text_buf)
            self.append_blocks(
                doc_token,
                [
                    {
                        "block_type": 5,
                        "heading3": {
                            "elements": [
                                {"text_run": {"content": f"—— 第 {slide_num} 页 ——"}}
                            ]
                        },
                    }
                ],
            )

    @staticmethod
    def _placeholder(msg: str):
        return {
            "block_type": 2,
            "text": {
                "elements": [
                    {
                        "text_run": {
                            "content": f"📎 {msg}",
                            "text_element_style": {"text_color": 4},
                        }
                    }
                ]
            },
        }
